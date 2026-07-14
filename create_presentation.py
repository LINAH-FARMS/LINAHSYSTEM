import urllib.request, urllib.error, json, os, io, math
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

SUPABASE_URL = 'https://cwqghiqykohefaggedjl.supabase.co'
SUPABASE_KEY = 'eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJpc3MiOiJzdXBhYmFzZSIsInJlZiI6ImN3cWdoaXF5a29oZWZhZ2dlZGpsIiwicm9sZSI6ImFub24iLCJpYXQiOjE3ODEwMjUyMjEsImV4cCI6MjA5NjYwMTIyMX0.3a3hRcNdmYQCtjYjBroAT6df1T_7oz-XWUeD3wagYw8'
ENDPOINT = SUPABASE_URL + '/rest/v1/sync_data'

OUTPUT = r'C:\Users\Salem Magdy\Desktop\LINAHSYSTEM.pptx'

plt.rcParams['font.family'] = 'sans-serif'
_arabic_fonts = [f.name for f in fm.fontManager.ttflist if any(n in f.name for n in ['Arabic','Arial','Times New Roman','Calibri','Tahoma'])]
if _arabic_fonts:
    plt.rcParams['font.sans-serif'] = _arabic_fonts
plt.rcParams['font.size'] = 13
plt.rcParams['axes.unicode_minus'] = False

# Colors
DARK_GREEN = RGBColor(0x1B, 0x5E, 0x20)
MED_GREEN = RGBColor(0x2E, 0x7D, 0x32)
LIGHT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
GOLD = RGBColor(0xFF, 0x8F, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG = RGBColor(0x1A, 0x23, 0x2F)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
TEXT_DARK = RGBColor(0x33, 0x33, 0x33)
TEXT_GRAY = RGBColor(0x75, 0x75, 0x75)
ACCENT = RGBColor(0x00, 0x96, 0x88)
RED = RGBColor(0xC6, 0x28, 0x28)
BLUE = RGBColor(0x15, 0x65, 0xC0)

# =============== FETCH DATA ===============
def fetch_supabase():
    req = urllib.request.Request(ENDPOINT + '?id=eq.alldata&select=data')
    req.add_header('apikey', SUPABASE_KEY)
    req.add_header('Authorization', 'Bearer ' + SUPABASE_KEY)
    with urllib.request.urlopen(req, timeout=30) as resp:
        rows = json.loads(resp.read())
    if rows and rows[0] and rows[0].get('data'):
        d = rows[0]['data']
        if isinstance(d, str): d = json.loads(d)
        return d
    return {}

data = fetch_supabase()

employees = data.get('employees', [])
rooms = data.get('roomsCapacity', [])
vacations = data.get('vacations', [])
meal_waste = data.get('mealWaste', [])
septic = data.get('septicRecords', [])
water_stations = data.get('waterStations', [])
hospitality = data.get('hospitalities', [])
maintenance = data.get('maintenanceRecords', [])
periodic_maint = data.get('periodicMaintenance', [])
meal_logs = data.get('mealLogs', [])
contractors = data.get('contractors', [])
inventory = data.get('inventoryItems', [])
evaluations = data.get('evaluations', [])
bakery_prod = data.get('bakeryProductions', [])
bakery_ctr = data.get('bakeryContractorSupplies', [])

# Stats
total_emp = len(employees)
p_count = sum(1 for e in employees if e.get('status') == 'P')
v_count = sum(1 for e in employees if e.get('status') == 'V')
x_count = sum(1 for e in employees if e.get('status') == 'X')
total_beds = sum(int(r.get('beds', 0)) for r in rooms)
occupied_beds = sum(1 for e in employees if (e.get('status') in ('P','V')) and e.get('room'))
vacant_beds = max(0, total_beds - occupied_beds)
emp_with_room = sum(1 for e in employees if e.get('room'))
emp_no_room = total_emp - emp_with_room

# Dep ts
depts = {}
for e in employees:
    d = e.get('dept', e.get('department', 'غير محدد'))
    depts[d] = depts.get(d, 0) + 1
sorted_depts = sorted(depts.items(), key=lambda x: -x[1])

# Waste stats
total_waste_meals = len(meal_waste)
total_prepared = sum(w.get('totalPrepared', 0) for w in meal_waste)
total_waste_val = sum((w.get('wasteEng', 0) or 0) + (w.get('wasteWrk', 0) or 0) for w in meal_waste)

print(f"Employees: {total_emp}, P={p_count}, V={v_count}, X={x_count}")
print(f"Beds: {total_beds}, Occupied: {occupied_beds}, Vacant: {vacant_beds}")
print(f"Depts: {len(sorted_depts)}")
print(f"Waste records: {total_waste_meals}")

# =============== CHARTS ===============
def save_chart(fname):
    path = os.path.join(os.path.dirname(OUTPUT), fname)
    plt.tight_layout(pad=1.5)
    plt.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close()
    return path

chart_files = []

# ---- Chart 1: Employee Status Pie ----
fig, ax = plt.subplots(figsize=(5.5, 5.5))
labels = ['موجود (P)', 'إجازة (V)', 'منتهي (X)']
sizes = [p_count, v_count, x_count]
colors = ['#2e7d32', '#1565c0', '#c62828']
explode = (0.03, 0.03, 0.03)
wedges, texts, autotexts = ax.pie(sizes, explode=explode, labels=labels, colors=colors,
    autopct='%1.1f%%', startangle=90, pctdistance=0.6,
    textprops={'fontsize': 12, 'fontweight': 'bold'})
for at in autotexts: at.set_color('white'); at.set_fontweight('bold')
ax.set_title('حالة الموظفين', fontsize=16, fontweight='bold', pad=15)
chart_files.append(save_chart('chart1_emp_status.png'))

# ---- Chart 2: Beds Overview ----
fig, ax = plt.subplots(figsize=(5.5, 4.5))
bars = ax.bar(['الأسرة المشغولة', 'الأسرة الخالية', 'إجمالي السعة'], 
              [occupied_beds, vacant_beds, total_beds],
              color=['#2e7d32', '#c62828', '#1565c0'], width=0.5)
for bar in bars:
    ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + max(0.5, bar.get_height()*0.02),
            str(int(bar.get_height())), ha='center', va='bottom', fontweight='bold', fontsize=13)
ax.set_ylabel('عدد الأسرة', fontsize=12)
ax.set_title('إشغال السكن', fontsize=16, fontweight='bold')
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
ax.set_ylim(0, total_beds * 1.2)
chart_files.append(save_chart('chart2_beds.png'))

# ---- Chart 3: Top Departments ----
top_depts = sorted_depts[:8]
fig, ax = plt.subplots(figsize=(6.5, 4.5))
dept_names = [d[0][:15] for d in top_depts]
dept_counts = [d[1] for d in top_depts]
bars = ax.barh(range(len(dept_names)), dept_counts, color='#009688', height=0.6)
ax.set_yticks(range(len(dept_names)))
ax.set_yticklabels(dept_names, fontsize=10)
ax.set_xlabel('عدد الموظفين', fontsize=12)
ax.set_title('أكبر 8 إدارات', fontsize=16, fontweight='bold')
for i, (bar, v) in enumerate(zip(bars, dept_counts)):
    ax.text(v + 0.3, i, str(v), va='center', fontweight='bold', fontsize=10)
ax.invert_yaxis()
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
chart_files.append(save_chart('chart3_depts.png'))

# ---- Chart 4: Employees with/without room ----
fig, ax = plt.subplots(figsize=(5, 4))
bars = ax.bar(['بسكن', 'بدون سكن'], [emp_with_room, emp_no_room],
              color=['#2e7d32', '#c62828'], width=0.5)
for bar in bars:
    ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1,
            str(int(bar.get_height())), ha='center', va='bottom', fontweight='bold', fontsize=13)
ax.set_title('توزيع السكن', fontsize=16, fontweight='bold')
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
ym = max(emp_with_room, emp_no_room) * 1.2
ax.set_ylim(0, ym if ym > 5 else 10)
chart_files.append(save_chart('chart4_housing.png'))

# ---- Chart 5: Meal Waste summary ----
waste_per_month = {}
for w in meal_waste:
    d = w.get('date', '')
    if d and len(d) >= 7:
        m = d[:7]
        waste_per_month[m] = waste_per_month.get(m, 0) + 1
if waste_per_month:
    sorted_months = sorted(waste_per_month.items())
    months = [m[0][:7] for m in sorted_months]
    counts = [m[1] for m in sorted_months]
    fig, ax = plt.subplots(figsize=(6.5, 4))
    ax.plot(range(len(months)), counts, marker='o', color='#e65100', linewidth=2.5, markersize=7)
    ax.fill_between(range(len(months)), counts, alpha=0.15, color='#e65100')
    ax.set_xticks(range(len(months)))
    ax.set_xticklabels([m[0][:7] for m in sorted_months], fontsize=9, rotation=30)
    ax.set_ylabel('عدد مرات تسجيل الهدر', fontsize=11)
    ax.set_title('تسجيلات هدر الوجبات', fontsize=16, fontweight='bold')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
else:
    fig, ax = plt.subplots(figsize=(5, 3))
    ax.text(0.5, 0.5, 'لا توجد بيانات', ha='center', va='center', fontsize=14)
chart_files.append(save_chart('chart5_waste.png'))

# ---- Chart 6: Total Prepared vs Waste ----
if meal_waste:
    total_waste_qty = sum((w.get('wasteEng', 0) or 0) + (w.get('wasteWrk', 0) or 0) for w in meal_waste)
    fig, ax = plt.subplots(figsize=(5.5, 4.5))
    cats = ['الإجمالي المُعد', 'إجمالي الهدر']
    vals = [total_prepared, total_waste_qty]
    bars = ax.bar(cats, vals, color=['#2e7d32', '#c62828'], width=0.4)
    for bar in bars:
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + max(1, bar.get_height()*0.02),
                f'{int(bar.get_height()):,}', ha='center', va='bottom', fontweight='bold', fontsize=12)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.set_title('إجمالي الوجبات مقابل الهدر', fontsize=14, fontweight='bold')
    chart_files.append(save_chart('chart6_waste_vs.png'))

# =============== BUILD PRESENTATION ===============
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_slide(layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=TEXT_DARK, align=PP_ALIGN.RIGHT, font_name='Arial'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rounded_rect(slide, left, top, width, height, color, opacity=None):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ===== SLIDE 1: TITLE =====
s1 = add_slide()
add_bg(s1, DARK_BG)
add_textbox(s1, 0.5, 1.2, 12, 1.5, 'LINAHSYSTEM', 50, True, GOLD, PP_ALIGN.CENTER)
add_textbox(s1, 0.5, 2.6, 12, 1.2, 'نظام إدارة مزرعة لينه فارمز', 40, True, WHITE, PP_ALIGN.CENTER)
add_textbox(s1, 0.5, 3.8, 12, 0.8, 'منصة متكاملة لإدارة الموظفين، السكن، المخزون، الوجبات، الصيانة، والمقاولين', 18, False, RGBColor(0xB0, 0xBE, 0xC5), PP_ALIGN.CENTER)
add_textbox(s1, 0.5, 5.0, 12, 0.5, 'مزودة بتقنية المزامنة السحابية (Supabase) والتشغيل على GitHub Pages', 16, False, RGBColor(0x78, 0x90, 0x9C), PP_ALIGN.CENTER)
# Decorative line
shape = s1.shapes.add_shape(1, Inches(4), Inches(4.8), Inches(5.333), Inches(0.04))
shape.fill.solid(); shape.fill.fore_color.rgb = GOLD; shape.line.fill.background()

# ===== SLIDE 2: SYSTEM OVERVIEW =====
s2 = add_slide()
add_bg(s2, WHITE)
add_rounded_rect(s2, 0, 0, 13.333, 1.2, DARK_GREEN)
add_textbox(s2, 0.5, 0.15, 12, 0.9, 'نظرة عامة على النظام', 34, True, WHITE, PP_ALIGN.RIGHT)

features = [
    ('👥 إدارة القوة العاملة', f'إدارة بيانات {total_emp} موظف مع تتبع الحالة (موجود، إجازة، منتهي)'),
    ('🏠 إدارة السكن', f'{total_beds} سرير في {len(rooms)} غرفة موزعة على قطاعات المزرعة'),
    ('📋 تقارير يومية وشهرية', 'تقارير شاملة: إجازات، ضيافة، صيانة، بيارات، حركة مخزن، وجبات'),
    ('☁️ مزامنة سحابية', 'مزامنة تلقائية مع Supabase — كل التغييرات تُحفظ فوراً'),
    ('💬 بوت واتساب', 'بوت للإبلاغ عن الأعطال (WhatsApp Bot)'),
    ('📱 تقارير على الموبايل', 'واجهة متجاوبة للهواتف (PWA-ready)'),
]
for i, (title, desc) in enumerate(features):
    y = 1.5 + i * 0.95
    add_rounded_rect(s2, 0.5, y, 12.3, 0.85, RGBColor(0xE8, 0xF5, 0xE9))
    add_textbox(s2, 0.7, y + 0.08, 11.5, 0.4, title, 18, True, DARK_GREEN)
    add_textbox(s2, 0.7, y + 0.42, 11.5, 0.35, desc, 14, False, TEXT_GRAY)

# Stats bar at bottom
add_rounded_rect(s2, 0.3, 7.0, 12.7, 0.4, DARK_BG)
stats_text = f'إجمالي الموظفين: {total_emp}  |  الأسرة: {total_beds}  |  المشغول: {occupied_beds}  |  الخالي: {vacant_beds}  |  الإدارات: {len(sorted_depts)}'
add_textbox(s2, 0.5, 7.02, 12.5, 0.35, stats_text, 12, True, WHITE, PP_ALIGN.CENTER)

# ===== SLIDE 3: EMPLOYEE STATS =====
s3 = add_slide()
add_bg(s3, WHITE)
add_rounded_rect(s3, 0, 0, 13.333, 1.2, MED_GREEN)
add_textbox(s3, 0.5, 0.15, 12, 0.9, 'إحصائيات القوة العاملة', 34, True, WHITE, PP_ALIGN.RIGHT)

s3.shapes.add_picture(chart_files[0], Inches(0.3), Inches(1.4), Inches(4.5), Inches(4.5))
s3.shapes.add_picture(chart_files[2], Inches(5.1), Inches(1.4), Inches(5.8), Inches(4.0))

# KPI boxes
kpis = [
    (0.3, 6.0, '👥 إجمالي', str(total_emp), DARK_GREEN),
    (3.5, 6.0, '✅ موجود', str(p_count), MED_GREEN),
    (6.7, 6.0, '✈️ إجازة', str(v_count), BLUE),
    (9.9, 6.0, '❌ منتهي', str(x_count), RED),
]
for x, y, lbl, val, c in kpis:
    add_rounded_rect(s3, x, y, 3.0, 1.2, c)
    add_textbox(s3, x, y + 0.1, 3.0, 0.45, lbl, 14, True, WHITE, PP_ALIGN.CENTER)
    add_textbox(s3, x, y + 0.5, 3.0, 0.55, val, 28, True, WHITE, PP_ALIGN.CENTER)

# ===== SLIDE 4: HOUSING =====
s4 = add_slide()
add_bg(s4, WHITE)
add_rounded_rect(s4, 0, 0, 13.333, 1.2, LIGHT_GREEN)
add_textbox(s4, 0.5, 0.15, 12, 0.9, 'إدارة السكن', 34, True, WHITE, PP_ALIGN.RIGHT)

s4.shapes.add_picture(chart_files[1], Inches(0.3), Inches(1.4), Inches(5.5), Inches(4.5))
s4.shapes.add_picture(chart_files[3], Inches(6.2), Inches(1.4), Inches(4.5), Inches(4.5))

# Room info
add_rounded_rect(s4, 0.3, 6.0, 12.7, 1.2, RGBColor(0xE8, 0xF5, 0xE9))
room_text = f'إجمالي الأسرة: {total_beds} سرير  |  مشغولة: {occupied_beds} سرير ({occupied_beds/total_beds*100:.0f}%)  |  خالية: {vacant_beds} سرير ({vacant_beds/total_beds*100:.0f}%)'
add_textbox(s4, 0.5, 6.1, 12, 0.4, room_text, 16, True, DARK_GREEN, PP_ALIGN.CENTER)
room_text2 = f'عدد الغرف: {len(rooms)} غرفة  |  موظفون بسكن: {emp_with_room}  |  موظفون بدون سكن: {emp_no_room}'
add_textbox(s4, 0.5, 6.5, 12, 0.4, room_text2, 14, False, TEXT_GRAY, PP_ALIGN.CENTER)

# ===== SLIDE 5: WASTE MANAGEMENT =====
s5 = add_slide()
add_bg(s5, WHITE)
add_rounded_rect(s5, 0, 0, 13.333, 1.2, GOLD)
add_textbox(s5, 0.5, 0.15, 12, 0.9, 'إدارة هدر الوجبات', 34, True, WHITE, PP_ALIGN.RIGHT)

total_waste_days = len(set(w.get('date', '') for w in meal_waste))
avg_daily = round(total_prepared / total_waste_days, 1) if total_waste_days else 0

s5.shapes.add_picture(chart_files[4], Inches(0.3), Inches(1.4), Inches(6.0), Inches(3.8))
if total_waste_meals:
    s5.shapes.add_picture(chart_files[5], Inches(6.8), Inches(1.4), Inches(5.0), Inches(3.7))

kpis_waste = [
    (0.3, 5.5, '📝 عدد التسجيلات', str(total_waste_meals), DARK_GREEN),
    (3.5, 5.5, '📆 أيام التسجيل', str(total_waste_days), BLUE),
    (6.7, 5.5, '🍽️ إجمالي المُعد', f'{int(total_prepared):,} كجم', GOLD),
    (9.9, 5.5, '🗑️ إجمالي الهدر', f'{int(total_waste_val):,}', RED),
]
for x, y, lbl, val, c in kpis_waste:
    add_rounded_rect(s5, x, y, 3.0, 1.1, c)
    add_textbox(s5, x, y + 0.08, 3.0, 0.4, lbl, 13, True, WHITE, PP_ALIGN.CENTER)
    add_textbox(s5, x, y + 0.45, 3.0, 0.5, val, 24, True, WHITE, PP_ALIGN.CENTER)

add_rounded_rect(s5, 0.3, 6.8, 12.7, 0.5, RGBColor(0xE8, 0xF5, 0xE9))
pct_waste = (total_waste_val / total_prepared * 100) if total_prepared else 0
add_textbox(s5, 0.5, 6.82, 12, 0.4, f'نسبة الهدر: {pct_waste:.1f}%  |  متوسط الإعداد اليومي: {avg_daily} كجم  |  تم تتبع التكلفة عبر قائمة المكونات الرئيسية', 14, True, DARK_GREEN, PP_ALIGN.CENTER)

# ===== SLIDE 6: REPORTS =====
s6 = add_slide()
add_bg(s6, WHITE)
add_rounded_rect(s6, 0, 0, 13.333, 1.2, BLUE)
add_textbox(s6, 0.5, 0.15, 12, 0.9, 'التقارير المتاحة', 34, True, WHITE, PP_ALIGN.RIGHT)

reports_list = [
    ('📊 تقارير القوة العاملة', ['توزيع الموظفين حسب الإدارة', 'حالة الحضور (موجود / إجازة / منتهي)', 'كشوف العهدة للموظفين', 'طباعة شاملة']),
    ('🏠 تقارير السكن', ['إشغال الأسرة في كل قطاع', 'توزيع الموظفين على الغرف', 'الموظفون بدون سكن', 'طباعة كشف السكن']),
    ('🍽️ تقارير الوجبات والهدر', ['تسجيلات هدر الوجبات اليومية', 'التكلفة لكل فرد', 'إجمالي الهدر مقابل المُعد', 'تصدير Excel']),
    ('🚛 تقارير البيارات', ['تسجيل نقلات البيارات', 'طباعة تقارير البيارات', 'تصدير Excel']),
    ('🛠️ تقارير الصيانة', ['سجلات الصيانة', 'الصيانة الدورية', 'تقارير الصيانة الشاملة', 'تصدير Excel']),
    ('🏪 تقارير المخزون', ['حركة المخزون', 'المخزونات', 'إذونات الصرف', 'أرشيف المخزون']),
    ('🍞 تقارير المخبز', ['إنتاج المخبز', 'توريد الخبز للمقاولين', 'المخزون والمكونات', 'فواتير المخبز']),
    ('📱 بلاغات الأعطال', ['بلاغات واتساب', 'تتبع حالة البلاغ', 'تقارير الأعطال']),
]
for i, (hdr, items) in enumerate(reports_list):
    col = i % 4
    row = i // 4
    x = 0.4 + col * 3.2
    y = 1.5 + row * 2.8
    add_rounded_rect(s6, x, y, 3.05, 2.6, RGBColor(0xE3, 0xF2, 0xFD))
    add_textbox(s6, x + 0.1, y + 0.08, 2.85, 0.4, hdr, 14, True, BLUE)
    for j, item in enumerate(items):
        add_textbox(s6, x + 0.2, y + 0.5 + j * 0.45, 2.7, 0.4, '• ' + item, 11, False, TEXT_DARK)

# ===== SLIDE 7: SYSTEM MODULES =====
s7 = add_slide()
add_bg(s7, WHITE)
add_rounded_rect(s7, 0, 0, 13.333, 1.2, RGBColor(0x00, 0x96, 0x88))
add_textbox(s7, 0.5, 0.15, 12, 0.9, 'وحدات النظام', 34, True, WHITE, PP_ALIGN.RIGHT)

modules = [
    ('الإدارة المرنة', 'إضافة وتعديل القطاعات، الغرف، البيارات، الإدارات، المسميات ديناميكياً'),
    ('المخزن', 'إدارة المخزون، إذونات الصرف، الجرد، والأرشيف'),
    ('الضيافة', 'تسجيل الضيوف وإدارة الضيافة'),
    ('الصيانة', 'سجلات الصيانة الدورية والطارئة'),
    ('المقاولين', 'إدارة عقود وبيانات المقاولين'),
    ('البيارات', 'تسجيل نقلات البيارات',
    ),
    ('الشاي والسكر', 'تسجيل صرف الشاي والسكر'),
    ('محطات المياه', 'تسجيل صيانة محطات المياه ورفع المستندات'),
    ('المخبز', 'إنتاج المخبز، المكونات، التوريد للمقاولين'),
    ('التقييمات', 'تقييم أداء الموظفين'),
    ('المركز المالي', 'الميزانية والمعاملات المالية'),
    ('سجل التعديلات', 'تسجيل جميع التعديلات على النظام'),
    ('المساعد الذكي', 'مساعد AI للاستعلامات'),
    ('تبادل البيانات', 'مزامنة مع السحابة والنسخ الاحتياطي'),
]

for i, (mod_name, mod_desc) in enumerate(modules):
    col = i % 4
    row = i // 4
    x = 0.3 + col * 3.2
    y = 1.4 + row * 1.45
    c = [RGBColor(0xE0, 0xF2, 0xF1), RGBColor(0xE8, 0xF5, 0xE9), RGBColor(0xE3, 0xF2, 0xFD), RGBColor(0xFD, 0xE7, 0xF6)][col % 4]
    add_rounded_rect(s7, x, y, 3.1, 1.3, c)
    add_textbox(s7, x + 0.1, y + 0.08, 2.9, 0.4, mod_name, 14, True, ACCENT)
    add_textbox(s7, x + 0.1, y + 0.48, 2.9, 0.7, mod_desc, 10, False, TEXT_GRAY)

# ===== SLIDE 8: BENEFITS =====
s8 = add_slide()
add_bg(s8, WHITE)
add_rounded_rect(s8, 0, 0, 13.333, 1.2, DARK_GREEN)
add_textbox(s8, 0.5, 0.15, 12, 0.9, 'المكتسبات والفوائد', 34, True, WHITE, PP_ALIGN.RIGHT)

benefits = [
    ('📈 تحسين إدارة القوى العاملة', 'تتبع دقيق لحضور وانصراف الموظفين وإجازاتهم'),
    ('🏠 إدارة السكن بكفاءة', f'توزيع {total_emp} موظف على {total_beds} سرير وتحديد الشواغر فوراً'),
    ('📊 تقارير لحظية دقيقة', 'تقارير شاملة بدون تدخل يدوي — وفر وقت وجهد'),
    ('☁️ مزامنة سحابية آمنة', 'البيانات محفوظة في Supabase — لا خوف من الفقدان'),
    ('🌐 تشغيل عبر الإنترنت', f'النظام متاح على GitHub Pages — دخول من أي جهاز'),
    ('🗑️ تقليل هدر الوجبات', f'تتبع {int(total_prepared):,} كجم من الوجبات بنسبة هدر {pct_waste:.1f}% — لترشيد الاستهلاك'),
    ('🔄 تكامل مع واتساب', 'استقبال بلاغات الأعطال عبر بوت واتساب آلياً'),
    ('💾 نسخ احتياطي تلقائي', 'نسخ احتياطي يومي في IndexedDB + سحابة Supabase'),
]
for i, (title, desc) in enumerate(benefits):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.2
    y = 1.5 + row * 1.35
    add_rounded_rect(s8, x, y, 5.9, 1.2, RGBColor(0xE8, 0xF5, 0xE9))
    add_textbox(s8, x + 0.2, y + 0.08, 5.5, 0.4, title, 16, True, DARK_GREEN)
    add_textbox(s8, x + 0.2, y + 0.5, 5.5, 0.5, desc, 13, False, TEXT_GRAY)

# ===== SLIDE 9: TECH STATS =====
s9 = add_slide()
add_bg(s9, WHITE)
add_rounded_rect(s9, 0, 0, 13.333, 1.2, DARK_BG)
add_textbox(s9, 0.5, 0.15, 12, 0.9, 'إحصائيات فنية', 34, True, GOLD, PP_ALIGN.RIGHT)

tech_stats = [
    ('الموظفين', str(total_emp), '👥'),
    ('الغرف', str(len(rooms)), '🏠'),
    ('الأسرة', str(total_beds), '🛏️'),
    ('الإجازات', str(len(vacations)), '✈️'),
    ('تسجيلات الضيافة', str(len(hospitality)), '🤝'),
    ('سجلات الصيانة', str(len(maintenance)), '🛠️'),
    ('الصيانة الدورية', str(len(periodic_maint)), '🔄'),
    ('نقلات البيارات', str(len(septic)), '🚛'),
    ('هدر الوجبات', str(total_waste_meals), '🗑️'),
    ('المقاولين', str(len(contractors)), '👷'),
    ('أصناف المخزون', str(len(inventory)), '📦'),
    ('محطات المياه', str(len(water_stations)), '💧'),
    ('تقييمات', str(len(evaluations)), '⭐'),
    ('إنتاج المخبز', str(len(bakery_prod)), '🍞'),
    ('توريد خبز', str(len(bakery_ctr)), '🚚'),
    ('تقارير الوجبات', str(len(meal_logs)), '🍽️'),
]

for i, (lbl, val, icon) in enumerate(tech_stats):
    col = i % 4
    row = i // 4
    x = 0.3 + col * 3.2
    y = 1.5 + row * 1.4
    add_rounded_rect(s9, x, y, 3.05, 1.25, RGBColor(0x26, 0x32, 0x38))
    add_textbox(s9, x + 0.1, y + 0.1, 2.85, 0.35, f'{icon} {lbl}', 12, True, GOLD, PP_ALIGN.CENTER)
    add_textbox(s9, x + 0.1, y + 0.5, 2.85, 0.5, val, 26, True, WHITE, PP_ALIGN.CENTER)

# ===== SLIDE 10: CLOSING =====
s10 = add_slide()
add_bg(s10, DARK_BG)
add_textbox(s10, 0.5, 1.5, 12, 1, 'شكراً لمتابعتكم', 44, True, GOLD, PP_ALIGN.CENTER)
add_textbox(s10, 0.5, 2.8, 12, 0.7, 'نظام إدارة مزرعة لينه فارمز', 28, True, WHITE, PP_ALIGN.CENTER)
add_textbox(s10, 0.5, 3.6, 12, 0.7, 'LINAHSYSTEM', 24, True, RGBColor(0xB0, 0xBE, 0xC5), PP_ALIGN.CENTER)
shape = s10.shapes.add_shape(1, Inches(4.5), Inches(4.3), Inches(4.333), Inches(0.04))
shape.fill.solid(); shape.fill.fore_color.rgb = GOLD; shape.line.fill.background()
add_textbox(s10, 0.5, 4.6, 12, 0.5, 'للتواصل: إدارة مزرعة لينه فارمز', 16, False, RGBColor(0x78, 0x90, 0x9C), PP_ALIGN.CENTER)
add_textbox(s10, 0.5, 5.2, 12, 0.5, 'https://linah-farms.github.io/LINAHSYSTEM', 14, False, RGBColor(0x54, 0x6E, 0x7A), PP_ALIGN.CENTER)

# =============== SAVE ===============
prs.save(OUTPUT)
print(f'\nDone. Presentation saved to: {OUTPUT}')
print(f'Slides: {len(prs.slides)}')
